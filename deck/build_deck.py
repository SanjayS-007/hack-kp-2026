"""
AEGIS pitch deck generator — HACK-KP 2026.
Builds a 16:9 dark-themed PPTX using only native python-pptx shapes.

Run (Windows PowerShell 5.1):
    $env:PYTHONPATH="C:\\Users\\2504690\\hack-kp-2026\\deck\\libs"
    & "C:\\Program Files\\Microsoft SDKs\\Azure\\CLI2\\python.exe" `
        "C:\\Users\\2504690\\hack-kp-2026\\deck\\build_deck.py"
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.oxml.ns import qn

# ----------------------------------------------------------------------------
# Design system
# ----------------------------------------------------------------------------
NAVY   = RGBColor(0x0A, 0x14, 0x26)   # base background
PANEL  = RGBColor(0x13, 0x21, 0x3C)   # card panel
PANEL2 = RGBColor(0x1B, 0x2C, 0x4C)   # lighter panel
STROKE = RGBColor(0x27, 0x3A, 0x5E)   # subtle border
CYAN   = RGBColor(0x22, 0xD3, 0xEE)   # accent
CYAN_D = RGBColor(0x0E, 0x74, 0x90)   # deep cyan
WHITE  = RGBColor(0xF8, 0xFA, 0xFC)
GRAY   = RGBColor(0x94, 0xA3, 0xB8)   # muted text
LGRAY  = RGBColor(0xCB, 0xD5, 0xE1)
AMBER  = RGBColor(0xF5, 0x9E, 0x0B)
RED    = RGBColor(0xEF, 0x44, 0x44)
GREEN  = RGBColor(0x34, 0xD3, 0x99)
VIOLET = RGBColor(0xA7, 0x8B, 0xFA)

FONT = "Segoe UI"
FONT_L = "Segoe UI Light"
FONT_SB = "Segoe UI Semibold"

EMU_W = Inches(13.333)
EMU_H = Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ----------------------------------------------------------------------------
# Helpers
# ----------------------------------------------------------------------------
def slide():
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, EMU_W, EMU_H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = NAVY
    bg.line.fill.background()
    bg.shadow.inherit = False
    return s


def _set_line(shape, color, width=1.0):
    if color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = color
        shape.line.width = Pt(width)


def rect(s, x, y, w, h, fill=PANEL, line=None, lw=1.0, rounded=True, radius=0.08):
    shp_type = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    shp = s.shapes.add_shape(shp_type, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    _set_line(shp, line, lw)
    shp.shadow.inherit = False
    if rounded:
        try:
            shp.adjustments[0] = radius
        except Exception:
            pass
    return shp


def oval(s, x, y, w, h, fill=None, line=CYAN, lw=1.5):
    shp = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(w), Inches(h))
    if fill is None:
        shp.fill.background()
    else:
        shp.fill.solid()
        shp.fill.fore_color.rgb = fill
    _set_line(shp, line, lw)
    shp.shadow.inherit = False
    return shp


def line(s, x1, y1, x2, y2, color=CYAN, w=1.5, dash=None):
    ln = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    ln.line.color.rgb = color
    ln.line.width = Pt(w)
    ln.shadow.inherit = False
    if dash:
        d = ln.line._get_or_add_ln()
        pd = d.makeelement(qn('a:prstDash'), {'val': dash})
        d.append(pd)
    return ln


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
         space_after=4, line_spacing=1.0, wrap=True):
    """runs: list of paragraphs; each paragraph is list of (txt,size,color,bold,font)."""
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = wrap
    tf.auto_size = MSO_AUTO_SIZE.NONE
    tf.vertical_anchor = anchor
    for m in (tf.margin_left, ):
        pass
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space_after)
        p.space_before = Pt(0)
        p.line_spacing = line_spacing
        for (txt, size, color, bold, fnt) in para:
            r = p.add_run()
            r.text = txt
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = fnt
    return tb


def R(txt, size, color=WHITE, bold=False, fnt=FONT):
    return (txt, size, color, bold, fnt)


def title_block(s, kicker, heading, num):
    line(s, 0.7, 1.02, 2.1, 1.02, CYAN, 3)
    text(s, 0.7, 0.55, 9.0, 0.4, [[R(kicker.upper(), 12.5, CYAN, True, FONT_SB)]])
    text(s, 0.7, 1.08, 11.5, 0.9, [[R(heading, 30, WHITE, True, FONT_SB)]])
    text(s, 12.15, 0.5, 0.9, 0.5, [[R(num, 12, GRAY, False, FONT)]], align=PP_ALIGN.RIGHT)


def chip(s, x, y, w, h, label, color, txtcolor=None):
    c = rect(s, x, y, w, h, fill=None, line=color, lw=1.25, radius=0.5)
    text(s, x, y + 0.01, w, h, [[R(label, 9.5, txtcolor or color, True, FONT_SB)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    return c


def footer(s, n):
    text(s, 0.7, 7.06, 8, 0.3,
         [[R("AEGIS-X  ·  HACK-KP 2026  ·  Confidential — Authorized Agencies Only", 8, GRAY, False, FONT)]])
    text(s, 12.2, 7.06, 0.9, 0.3, [[R(str(n), 8, GRAY, False, FONT)]], align=PP_ALIGN.RIGHT)


# ----------------------------------------------------------------------------
# Slide 1 — Title
# ----------------------------------------------------------------------------
def s1():
    s = slide()
    # decorative concentric shield rings on the right
    oval(s, 8.7, 0.9, 5.6, 5.6, fill=None, line=RGBColor(0x14,0x2A,0x44), lw=1.2)
    oval(s, 9.4, 1.6, 4.2, 4.2, fill=None, line=RGBColor(0x18,0x33,0x52), lw=1.2)
    oval(s, 10.05, 2.25, 2.9, 2.9, fill=PANEL, line=CYAN, lw=1.75)
    shield = s.shapes.add_shape(MSO_SHAPE.HEXAGON, Inches(10.65), Inches(2.9), Inches(1.7), Inches(1.6))
    shield.fill.solid(); shield.fill.fore_color.rgb = CYAN_D
    shield.line.color.rgb = CYAN; shield.line.width = Pt(1.5); shield.shadow.inherit = False
    text(s, 10.65, 3.35, 1.7, 0.7, [[R("AI", 26, WHITE, True, FONT_SB)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)

    line(s, 0.85, 2.35, 2.35, 2.35, CYAN, 3)
    text(s, 0.8, 2.5, 8.6, 1.6, [[R("AEGIS", 78, WHITE, True, FONT_SB),
                                   R("-X", 78, CYAN, True, FONT_SB)]])
    text(s, 0.85, 3.9, 8.6, 0.6,
         [[R("AI Evidence & Guardianship Intelligence System", 19, CYAN, True, FONT_SB)]])
    text(s, 0.85, 4.42, 8.6, 0.4,
         [[R("The sovereign intelligence fabric for child protection — provable at every step.", 12.5, GRAY, False, FONT)]])
    text(s, 0.85, 4.9, 8.4, 0.6,
         [[R("From terabytes to testimony — in hours, not months.", 17, LGRAY, False, FONT)]])
    # bottom band
    rect(s, 0.85, 5.9, 8.4, 0.9, fill=PANEL, line=STROKE, lw=1, radius=0.12)
    text(s, 1.1, 5.9, 8.0, 0.9,
         [[R("HACK-KP 2026", 13, WHITE, True, FONT_SB),
           R("   ·   Child-Protection Investigation Track", 13, GRAY, False, FONT)],
          [R("Air-gapped  ·  Court-admissible  ·  Investigator-wellbeing-first", 10.5, CYAN, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=3)
    footer(s, 1)


# ----------------------------------------------------------------------------
# Slide 2 — The Problem
# ----------------------------------------------------------------------------
def s2():
    s = slide()
    title_block(s, "The Problem", "Investigators are drowning — and hashing is blind", "01")
    text(s, 0.7, 1.95, 11.9, 0.5,
         [[R("Exploding evidence volume collides with manual review, human trauma, and content that legacy tools cannot see.",
             13.5, LGRAY, False, FONT)]])

    cards = [
        ("2.4 TB", "per seized case", "480k+ files across 3 devices in a single operation — impossible to review by hand.", AMBER),
        ("9 months", "typical review backlog", "Manual triage delays justice while victims remain in active harm.", RED),
        ("Secondary\nTrauma", "investigator STS / PTSD", "Direct exposure to abuse imagery drives burnout and severe workforce turnover.", VIOLET),
        ("Zero-day", "invisible to hashing", "PhotoDNA / Project VIC hashes fail on novel & AI-generated content.", CYAN),
    ]
    x = 0.7; w = 2.86; gap = 0.13
    for i, (big, sub, body, col) in enumerate(cards):
        cx = x + i * (w + gap)
        rect(s, cx, 2.75, w, 2.65, fill=PANEL, line=STROKE, lw=1, radius=0.06)
        rect(s, cx, 2.75, 0.09, 2.65, fill=col, line=None, rounded=False)
        text(s, cx + 0.28, 2.95, w - 0.45, 1.1,
             [[R(big, 27 if "\n" not in big else 21, col, True, FONT_SB)]], anchor=MSO_ANCHOR.TOP)
        text(s, cx + 0.28, 4.02, w - 0.45, 0.4,
             [[R(sub.upper(), 10, WHITE, True, FONT_SB)]])
        text(s, cx + 0.28, 4.42, w - 0.5, 0.95,
             [[R(body, 10.5, GRAY, False, FONT)]], line_spacing=1.05)

    # market band
    rect(s, 0.7, 5.7, 11.93, 1.05, fill=PANEL2, line=CYAN, lw=1.25, radius=0.1)
    text(s, 1.0, 5.7, 5.2, 1.05,
         [[R("Digital-forensics market", 11, GRAY, False, FONT)],
          [R("$12.94B → $22.81B", 24, CYAN, True, FONT_SB)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    line(s, 6.4, 5.95, 6.4, 6.5, STROKE, 1)
    text(s, 6.7, 5.7, 6.0, 1.05,
         [[R("by 2030  ·  12.0% CAGR", 13, WHITE, True, FONT_SB)],
          [R("The surge is driven by IoT endpoints, encrypted volumes, and sophisticated evasion — automation is no longer optional.",
             10, GRAY, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=2)
    footer(s, 2)


# ----------------------------------------------------------------------------
# Slide 3 — Challenge Fit (12 areas grid)
# ----------------------------------------------------------------------------
def s3():
    s = slide()
    title_block(s, "Challenge Fit", "12 innovation areas → 12 AEGIS-X modules", "02")
    text(s, 0.7, 1.9, 11.9, 0.4,
         [[R("Every required capability maps to a purpose-built module in one unified platform.", 13, LGRAY, False, FONT)]])

    items = [
        ("01", "Content Analysis", "Visual Triage Engine"),
        ("02", "Threat Identification", "Human-Centric Perception"),
        ("03", "Source Correlation", "Entity Graph"),
        ("04", "Contextual Extraction", "Context Miner"),
        ("05", "Activity Patterns", "Behavior Profiler"),
        ("06", "Metadata Mapping", "Provenance Mapper"),
        ("07", "Synthetic Detection", "DeepFake Shield"),
        ("08", "Timeline Reconstruction", "Chronos Engine"),
        ("09", "Intelligent Retrieval", "Ask AEGIS"),
        ("10", "Automated Reporting", "Court-Ready Reports"),
        ("11", "Risk Assessment", "Lead Prioritizer"),
        ("12", "Intelligence Fusion", "Fusion Center"),
    ]
    cols, rows = 4, 3
    x0, y0 = 0.7, 2.55
    w, h = 2.92, 1.32
    gx, gy = 0.13, 0.14
    for idx, (num, area, mod) in enumerate(items):
        r = idx // cols; c = idx % cols
        cx = x0 + c * (w + gx); cy = y0 + r * (h + gy)
        rect(s, cx, cy, w, h, fill=PANEL, line=STROKE, lw=1, radius=0.07)
        text(s, cx + 0.22, cy + 0.14, 1.0, 0.4, [[R(num, 15, CYAN, True, FONT_SB)]])
        text(s, cx + 0.22, cy + 0.54, w - 0.4, 0.4,
             [[R(area, 11.5, WHITE, True, FONT_SB)]])
        text(s, cx + 0.22, cy + 0.9, w - 0.4, 0.4,
             [[R("▸ " + mod, 10, CYAN, False, FONT)]])
    footer(s, 3)


# ----------------------------------------------------------------------------
# Slide 4 — Solution Overview (pipeline)
# ----------------------------------------------------------------------------
def s4():
    s = slide()
    title_block(s, "Solution Overview", "One console. From seizure to sworn testimony.", "03")
    text(s, 0.7, 1.95, 11.9, 0.7,
         [[R("AEGIS-X is an air-gapped AI investigation platform that turns terabytes of seized evidence into a "
             "prioritized, court-ready case file — while shielding investigators from traumatic content.",
             14, LGRAY, False, FONT)]], line_spacing=1.1)

    stages = [
        ("INGEST", "Seized devices,\nVICS normalization", CYAN),
        ("ANALYZE", "Vision · NLP · age &\npose · deepfake", CYAN),
        ("CORRELATE", "Entity graph &\ntimeline fusion", CYAN),
        ("PRIORITIZE", "Severity scoring &\nactive-abuse queue", AMBER),
        ("REPORT", "§63 certificate &\nSHA-256 manifest", GREEN),
    ]
    n = len(stages)
    x0 = 0.7; y = 3.35; w = 2.15; h = 1.85; gap = 0.28
    for i, (name, sub, col) in enumerate(stages):
        cx = x0 + i * (w + gap)
        rect(s, cx, y, w, h, fill=PANEL, line=col, lw=1.5, radius=0.08)
        oval(s, cx + w/2 - 0.28, y + 0.22, 0.56, 0.56, fill=None, line=col, lw=1.75)
        text(s, cx + w/2 - 0.28, y + 0.22, 0.56, 0.56,
             [[R(str(i+1), 17, col, True, FONT_SB)]], align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx, y + 0.92, w, 0.4, [[R(name, 14, WHITE, True, FONT_SB)]], align=PP_ALIGN.CENTER)
        text(s, cx, y + 1.3, w, 0.5, [[R(sub, 9.5, GRAY, False, FONT)]],
             align=PP_ALIGN.CENTER, line_spacing=1.0)
        if i < n - 1:
            ax = cx + w + 0.02
            line(s, ax, y + h/2, ax + gap - 0.04, y + h/2, col, 2.5)
            tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ax + gap - 0.14),
                                     Inches(y + h/2 - 0.08), Inches(0.16), Inches(0.16))
            tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = col
            tri.line.fill.background(); tri.shadow.inherit = False

    rect(s, 0.7, 5.65, 11.93, 1.05, fill=PANEL2, line=STROKE, lw=1, radius=0.1)
    outs = [("47 min", "auto-triage of 2.4 TB", CYAN),
            ("~91%", "less traumatic content seen", GREEN),
            ("100% local", "sovereign & air-gapped", VIOLET),
            ("Court-ready", "by design, not afterthought", AMBER)]
    ow = 11.93 / 4
    for i, (b, sub, col) in enumerate(outs):
        ox = 0.7 + i * ow
        if i > 0:
            line(s, ox, 5.85, ox, 6.5, STROKE, 1)
        text(s, ox + 0.25, 5.65, ow - 0.3, 1.05,
             [[R(b, 19, col, True, FONT_SB)], [R(sub, 10, GRAY, False, FONT)]],
             anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    footer(s, 4)


# ----------------------------------------------------------------------------
# Slide 5 — Architecture
# ----------------------------------------------------------------------------
def s5():
    s = slide()
    title_block(s, "Architecture", "Layered pipeline inside a Zero-Trust boundary", "04")

    # Zero-Trust boundary
    rect(s, 0.62, 2.05, 12.1, 4.55, fill=None, line=CYAN_D, lw=1.5, radius=0.03)
    chip(s, 0.95, 1.83, 3.7, 0.42, "AIR-GAPPED  ·  ZERO-TRUST PERIMETER", CYAN)
    chip(s, 8.55, 1.83, 4.1, 0.42, "CONFIDENTIAL-GPU ENCLAVE · ATTESTED ✓", VIOLET, VIOLET)

    layers = [
        ("EVIDENCE SOURCES", "Seized phones & PCs · CCTNS · CDR/IPDR · OSINT · blockchain", GRAY, STROKE),
        ("INGESTION  &  VICS NORMALIZATION", "Hashing (SHA-256) · EXIF/GPS · C2PA · Project VIC JSON schema", CYAN, CYAN_D),
        ("AI CORE · INVESTIGATION SWARM", "Orchestrator → Triage·Graph·Vision·Crypto·Chronos·Report → Verifier Agent", CYAN, CYAN),
        ("FUSION GRAPH", "One temporal graph: suspects · victims · hashes · wallets · IPs · timeline", VIOLET, VIOLET),
        ("INVESTIGATOR CONSOLE", "Blur-by-default triage · XAI overlays · Ask AEGIS · risk queue", WHITE, STROKE),
        ("COURT-READY OUTPUT", "BSA-2023 §63 · SHA-256 chain of custody · enclave attestation · ECS-gated", GREEN, GREEN),
    ]
    x = 1.0; y = 2.35; w = 11.35; h = 0.62; gap = 0.075
    for i, (name, sub, tcol, bcol) in enumerate(layers):
        cy = y + i * (h + gap)
        fill = PANEL2 if name in ("AI CORE · INVESTIGATION SWARM", "COURT-READY OUTPUT") else PANEL
        rect(s, x, cy, w, h, fill=fill, line=bcol, lw=1.25, radius=0.06)
        text(s, x + 0.3, cy, 5.0, h, [[R(name, 12, tcol, True, FONT_SB)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, x + 5.3, cy, w - 5.5, h, [[R(sub, 10, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)
        if name == "AI CORE · INVESTIGATION SWARM":
            chip(s, x + w - 1.72, cy + 0.14, 1.5, 0.34, "VERIFIER ✓", GREEN, GREEN)
        if i < len(layers) - 1:
            mx = x + w/2
            line(s, mx, cy + h, mx, cy + h + gap, bcol if bcol != STROKE else GRAY, 1.5)
    footer(s, 5)


# ----------------------------------------------------------------------------
# Slide 6 — AI Core Deep-Dive (4 quadrants)
# ----------------------------------------------------------------------------
def s6():
    s = slide()
    title_block(s, "AI Core Deep-Dive", "Four engines, explainable by design", "05")

    quads = [
        ("Hybrid Content Analysis", CYAN,
         ["ConvNeXt-Tiny + Swin Transformer fusion", "Grad-CAM XAI overlays on every flag",
          "Selective blur / semantic occlusion", "Beyond NSFW — grooming & context aware"]),
        ("Human-Centric Perception", AMBER,
         ["YOLO-Pose skeletal keypoints (BKPD)", "Apparent-age distribution via KDE",
          "Pose explicitness under obfuscation", "Flags likely-minor with statistical reliability"]),
        ("Entity Graph — TAGNN", VIOLET,
         ["Transformer-autoencoder graph NN", "Nodes: suspects · victims · hashes · wallets · IPs",
          "Address clustering on crypto peel-chains", "Cross-case linkage of hidden actors"]),
        ("Multimodal RAG — Ask AEGIS", GREEN,
         ["ChromaDB vectors + local Llama (on-prem)", "Citation-grounded, deterministic answers",
          "HERAM hallucination scoring → ECS", "Below-threshold output auto-excluded"]),
    ]
    x0, y0 = 0.7, 2.15
    w, h = 5.9, 2.28
    gx, gy = 0.13, 0.16
    for idx, (name, col, bullets) in enumerate(quads):
        r = idx // 2; c = idx % 2
        cx = x0 + c * (w + gx); cy = y0 + r * (h + gy)
        rect(s, cx, cy, w, h, fill=PANEL, line=STROKE, lw=1, radius=0.05)
        rect(s, cx, cy, w, 0.62, fill=PANEL2, line=None, radius=0.05)
        oval(s, cx + 0.28, cy + 0.16, 0.3, 0.3, fill=col, line=None)
        text(s, cx + 0.75, cy, w - 1, 0.62, [[R(name, 14.5, WHITE, True, FONT_SB)]], anchor=MSO_ANCHOR.MIDDLE)
        paras = [[R("•  ", 11, col, True, FONT), R(b, 11, LGRAY, False, FONT)] for b in bullets]
        text(s, cx + 0.35, cy + 0.75, w - 0.6, h - 0.85, paras, space_after=5, line_spacing=1.0)
    footer(s, 6)


# ----------------------------------------------------------------------------
# Slide 7 — Synthetic / Deepfake Shield
# ----------------------------------------------------------------------------
def s7():
    s = slide()
    title_block(s, "DeepFake Shield", "Catches what hash-matching can't", "06")
    text(s, 0.7, 1.9, 11.9, 0.4,
         [[R("Synthetic CSAM (Stable Diffusion, GANs, face-swap) creates brand-new pixels — hashes are useless. "
             "AEGIS runs a three-stream ensemble plus optical physics.", 12.5, LGRAY, False, FONT)]],
         line_spacing=1.05)

    streams = [
        ("1", "Global Texture", "DINOv2 backbone", "Macro-context & GAN fingerprint artifacts across the whole frame.", CYAN),
        ("2", "Facial Geometry", "Localized facial stream", "Anchors biometric landmarks; detects warped / swapped identity.", AMBER),
        ("3", "Semantic Fusion", "Frozen-CLIP backbone", "Flags logical impossibilities & scene inconsistencies.", VIOLET),
    ]
    x0 = 0.7; y = 2.75; w = 3.6; h = 2.05; gap = 0.28
    for i, (n, name, tech, body, col) in enumerate(streams):
        cx = x0 + i * (w + gap)
        rect(s, cx, y, w, h, fill=PANEL, line=col, lw=1.5, radius=0.06)
        oval(s, cx + 0.3, y + 0.28, 0.5, 0.5, fill=None, line=col, lw=1.75)
        text(s, cx + 0.3, y + 0.28, 0.5, 0.5, [[R(n, 16, col, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 0.95, y + 0.28, w - 1.1, 0.5, [[R(name, 14.5, WHITE, True, FONT_SB)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 0.3, y + 0.95, w - 0.55, 0.35, [[R(tech, 10.5, col, True, FONT_SB)]])
        text(s, cx + 0.3, y + 1.3, w - 0.55, 0.7, [[R(body, 10.5, GRAY, False, FONT)]], line_spacing=1.05)
        # arrow into fusion
        line(s, cx + w/2, y + h, cx + w/2, y + h + 0.22, col, 2)

    # Fusion verdict bar
    rect(s, 0.7, 5.35, 11.93, 1.35, fill=PANEL2, line=GREEN, lw=1.5, radius=0.08)
    text(s, 1.0, 5.5, 4.6, 1.1,
         [[R("ENSEMBLE  VERDICT", 11, GREEN, True, FONT_SB)],
          [R("Authentic  ·  Synthetic  ·  Manipulated", 13, WHITE, True, FONT_SB)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=3)
    line(s, 5.9, 5.55, 5.9, 6.5, STROKE, 1)
    text(s, 6.15, 5.5, 6.4, 1.15,
         [[R("+  Defocus-blur optics", 11.5, CYAN, True, FONT_SB),
           R("  — real lenses have physical depth-of-field; latent-space models don't.", 10.5, GRAY, False, FONT)],
          [R("+  A/V sync matrix", 11.5, CYAN, True, FONT_SB),
           R("  — lip / audio desynchronization exposes video fabrication.", 10.5, GRAY, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=6)
    footer(s, 7)


# ----------------------------------------------------------------------------
# Slide 8 — Court-Admissible by Design
# ----------------------------------------------------------------------------
def s8():
    s = slide()
    title_block(s, "Court-Admissible by Design", "Evidentiary integrity is built in, not bolted on", "07")

    cards = [
        ("SHA-256", "Chain of custody",
         "256-bit hash computed at the moment of extraction. One altered pixel breaks the hash — tampering is instantly provable.", CYAN),
        ("BSA 2023 §63", "Auto-certificates",
         "Two-part statutory certificate auto-populated with IMEI, MAC & serials — mandatory for admissible secondary e-evidence in India.", GREEN),
        ("VICS", "Interoperability",
         "Native Project VIC JSON: shares SourceID, hashes (MD5/SHA-256/PhotoDNA) & SAP grade with Cellebrite, Magnet, MSAB.", AMBER),
        ("ECS Gating", "Hallucination control",
         "HERAM scores every AI claim vs source text; output below the judicial threshold is flagged or excluded automatically.", VIOLET),
    ]
    x = 0.7; w = 2.86; gap = 0.13; y = 2.25; h = 3.05
    for i, (big, sub, body, col) in enumerate(cards):
        cx = x + i * (w + gap)
        rect(s, cx, y, w, h, fill=PANEL, line=STROKE, lw=1, radius=0.06)
        rect(s, cx, y, w, 0.12, fill=col, line=None, rounded=False)
        text(s, cx + 0.28, y + 0.35, w - 0.5, 0.6, [[R(big, 20, col, True, FONT_SB)]])
        text(s, cx + 0.28, y + 1.02, w - 0.5, 0.35, [[R(sub.upper(), 10, WHITE, True, FONT_SB)]])
        line(s, cx + 0.28, y + 1.42, cx + w - 0.28, y + 1.42, STROKE, 1)
        text(s, cx + 0.28, y + 1.55, w - 0.5, 1.4, [[R(body, 10.5, GRAY, False, FONT)]], line_spacing=1.1)

    rect(s, 0.7, 5.6, 11.93, 1.05, fill=PANEL2, line=CYAN, lw=1.25, radius=0.1)
    text(s, 1.0, 5.6, 11.4, 1.05,
         [[R("Every artifact carries a verifiable provenance trail — ", 12.5, LGRAY, False, FONT),
           R("structurally unassailable under cross-examination.", 12.5, CYAN, True, FONT_SB)]],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 8)


# ----------------------------------------------------------------------------
# Slide 9 — Demo Walkthrough (Operation Sentinel)
# ----------------------------------------------------------------------------
def s9():
    s = slide()
    title_block(s, "Demo Walkthrough", "\"Operation Sentinel\"  ·  Case #KP-2026-0417", "08")
    text(s, 0.7, 1.9, 11.9, 0.4,
         [[R("Fictional, 100% synthetic case data — no real identities or imagery.", 11.5, GRAY, False, FONT)]])

    steps = [
        ("1", "Ingest", "Ingest & Triage", "3 seized devices — 2.4 TB, 480k files — auto-triaged in 47 minutes.", CYAN),
        ("2", "Dashboard", "Risk Queue", "312 items flagged · 14 high-risk · 3 synthetic detected, severity-sorted.", AMBER),
        ("3", "Entity Graph", "Hidden Link", "Graph reveals a 3rd suspect via crypto peel-chain + shared hash.", VIOLET),
        ("4", "Timeline", "Chronos Engine", "Reconstructs grooming → production → distribution across devices.", CYAN),
        ("5", "Ask AEGIS", "RAG Query", "\"Show all comms near Riverside Park in March\" → cited answer.", GREEN),
        ("6", "Report", "§63 Export", "One-click court report with §63 certificate + SHA-256 manifest.", GREEN),
    ]
    cols = 3
    x0, y0 = 0.7, 2.5
    w, h = 3.87, 1.75
    gx, gy = 0.16, 0.18
    for idx, (n, page, name, body, col) in enumerate(steps):
        r = idx // cols; c = idx % cols
        cx = x0 + c * (w + gx); cy = y0 + r * (h + gy)
        rect(s, cx, cy, w, h, fill=PANEL, line=STROKE, lw=1, radius=0.06)
        # mini screenshot placeholder
        rect(s, cx + 0.22, cy + 0.22, 1.0, 1.3, fill=PANEL2, line=col, lw=1.25, radius=0.06)
        rect(s, cx + 0.34, cy + 0.36, 0.76, 0.16, fill=col, line=None, radius=0.3)
        for k in range(3):
            rect(s, cx + 0.34, cy + 0.62 + k*0.22, 0.76, 0.1, fill=STROKE, line=None, radius=0.4)
        text(s, cx + 0.22, cy + 1.28, 1.0, 0.25, [[R(page, 8, GRAY, False, FONT)]], align=PP_ALIGN.CENTER)
        # step content
        oval(s, cx + 1.4, cy + 0.2, 0.42, 0.42, fill=col, line=None)
        text(s, cx + 1.4, cy + 0.2, 0.42, 0.42, [[R(n, 14, NAVY, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 1.95, cy + 0.22, w - 2.1, 0.42, [[R(name, 12.5, WHITE, True, FONT_SB)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 1.42, cy + 0.72, w - 1.6, 0.9, [[R(body, 10, GRAY, False, FONT)]], line_spacing=1.05)
    footer(s, 9)


# ----------------------------------------------------------------------------
# Slide 10 — Impact
# ----------------------------------------------------------------------------
def s10():
    s = slide()
    title_block(s, "Impact", "Faster justice. Safer investigators.", "09")

    big = [
        ("~91%", "less traumatic\ncontent viewed by humans", "AI absorbs the psychological toxicity; examiners read clinical descriptions & blurred tiles.", GREEN),
        ("47 min", "vs  9 months", "Auto-triage of a 2.4 TB, 480k-file case — a step-change in throughput.", CYAN),
        ("72 hrs", "victim safeguarded", "Operation Sentinel outcome vs a typical multi-month timeline.", AMBER),
    ]
    x0 = 0.7; y = 2.35; w = 3.87; h = 2.9; gx = 0.16
    for i, (num, sub, body, col) in enumerate(big):
        cx = x0 + i * (w + gx)
        rect(s, cx, y, w, h, fill=PANEL, line=col, lw=1.5, radius=0.06)
        text(s, cx, y + 0.4, w, 1.1, [[R(num, 52, col, True, FONT_SB)]], align=PP_ALIGN.CENTER)
        text(s, cx + 0.3, y + 1.55, w - 0.6, 0.7, [[R(sub, 13.5, WHITE, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, line_spacing=1.0)
        line(s, cx + 0.5, y + 2.15, cx + w - 0.5, y + 2.15, STROKE, 1)
        text(s, cx + 0.3, y + 2.25, w - 0.6, 0.6, [[R(body, 10, GRAY, False, FONT)]],
             align=PP_ALIGN.CENTER, line_spacing=1.05)

    rect(s, 0.7, 5.55, 11.93, 1.1, fill=PANEL2, line=VIOLET, lw=1.5, radius=0.1)
    text(s, 1.0, 5.55, 11.4, 1.1,
         [[R("Investigator wellbeing is a first-class design goal — ", 13.5, LGRAY, False, FONT),
           R("sustainable teams clear backlogs, and cleared backlogs rescue children sooner.", 13.5, VIOLET, True, FONT_SB)]],
         anchor=MSO_ANCHOR.MIDDLE, line_spacing=1.05)
    footer(s, 10)


# ----------------------------------------------------------------------------
# Slide 11 — Differentiators (table)
# ----------------------------------------------------------------------------
def s11():
    s = slide()
    title_block(s, "Differentiators", "Why AEGIS-X, not the status quo", "10")

    from pptx.util import Inches as In
    rows = [
        ("Capability", "Hash tools\n(PhotoDNA / VIC)", "Forensic suites\n(Cellebrite / Griffeye)", "AEGIS-X"),
        ("Known CSAM match", "Yes", "Yes", "Yes"),
        ("Zero-day / novel content", "No", "Partial", "Yes — CNN+ViT semantics"),
        ("AI-generated / deepfake", "No", "Limited", "Yes — 3-stream + optics"),
        ("Cross-source fusion graph", "No", "Siloed", "Yes — one TAGNN graph"),
        ("Court-ready AI output", "N/A", "Manual", "Yes — §63 + ECS gating"),
        ("Investigator wellbeing", "No", "Basic blur", "Blur-by-default + XAI"),
    ]
    left, top = 0.7, 2.2
    tw, th = 11.93, 4.35
    gtable = s.shapes.add_table(len(rows), 4, In(left), In(top), In(tw), In(th)).table
    gtable.columns[0].width = In(3.35)
    gtable.columns[1].width = In(2.7)
    gtable.columns[2].width = In(2.9)
    gtable.columns[3].width = In(2.98)
    # disable banding style by setting first row + plain
    tbl = gtable._tbl
    # remove default style to allow custom fills
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gtable.cell(r, c)
            cell.margin_left = In(0.14); cell.margin_right = In(0.1)
            cell.margin_top = In(0.05); cell.margin_bottom = In(0.05)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2 if c < 3 else CYAN_D
            elif c == 3:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x0F, 0x2B, 0x33)
            elif c == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = val
            if r == 0:
                run.font.size = Pt(11.5); run.font.bold = True
                run.font.color.rgb = WHITE if c < 3 else NAVY
            elif c == 3:
                run.font.size = Pt(11); run.font.bold = True; run.font.color.rgb = CYAN
            elif c == 0:
                run.font.size = Pt(11.5); run.font.bold = True; run.font.color.rgb = WHITE
            else:
                run.font.size = Pt(11); run.font.color.rgb = GRAY
                if val in ("No", "N/A"):
                    run.font.color.rgb = RED
                elif val in ("Partial", "Limited", "Siloed", "Manual", "Basic blur"):
                    run.font.color.rgb = AMBER
                elif val == "Yes":
                    run.font.color.rgb = GREEN
            run.font.name = FONT
    footer(s, 11)


# ----------------------------------------------------------------------------
# Slide 12 — Roadmap & Ethics
# ----------------------------------------------------------------------------
def s12():
    s = slide()
    title_block(s, "Roadmap & Ethics", "Scale responsibly, with guardrails first", "15")

    # Roadmap timeline
    text(s, 0.7, 1.95, 6, 0.4, [[R("ROADMAP", 12, CYAN, True, FONT_SB)]])
    phases = [
        ("Phase 1", "Pilot + Edge Kit", "Kerala cyber-cell pilot on air-gapped hardware; AEGIS-X Edge NPU field kit seals the scene manifest.", CYAN),
        ("Phase 2", "Agency Integration", "National rollout via I4C / CCTNS + CDR/IPDR data lakes; multi-device fusion at scale.", AMBER),
        ("Phase 3", "Federated Guardian Network", "Cross-agency FL + secure aggregation + DP — every deployment improves detection; data never moves.", GREEN),
        ("Phase 4", "International", "INTERPOL / ICSE interoperability via VICS; sovereign-by-default across jurisdictions.", VIOLET),
    ]
    y = 2.42; ph = 1.06
    line(s, 1.05, y + 0.05, 1.05, y + 4*ph - 0.55, STROKE, 1.5)
    for i, (ph_n, name, body, col) in enumerate(phases):
        cy = y + i * ph
        oval(s, 0.85, cy + 0.05, 0.4, 0.4, fill=col, line=None)
        text(s, 0.85, cy + 0.05, 0.4, 0.4, [[R(str(i+1), 13, NAVY, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, 1.5, cy, 5.1, 0.4, [[R(ph_n + "  ·  ", 11, col, True, FONT_SB), R(name, 12.5, WHITE, True, FONT_SB)]])
        text(s, 1.5, cy + 0.4, 5.15, 0.6, [[R(body, 9.8, GRAY, False, FONT)]], line_spacing=1.0)

    # Ethics guardrails
    rect(s, 7.0, 2.35, 5.63, 4.3, fill=PANEL, line=VIOLET, lw=1.25, radius=0.05)
    text(s, 7.3, 2.55, 5, 0.4, [[R("ETHICS GUARDRAILS", 12, VIOLET, True, FONT_SB)]])
    guards = [
        ("Authorized-agency-only", "Deployed exclusively for vetted child-protection units."),
        ("Privacy by architecture", "Air-gapped, sovereign; no data ever leaves the agency."),
        ("Human-in-the-loop", "AI triages & recommends — humans decide and are accountable."),
        ("No training on illicit data", "Synthetic data & benign transfer learning; no CSAM ingested to train."),
        ("Bias safeguards", "Linear Sequential Unmasking + XAI to resist confirmation bias."),
    ]
    gy = 3.05
    for i, (t, b) in enumerate(guards):
        cy = gy + i * 0.72
        oval(s, 7.35, cy + 0.05, 0.22, 0.22, fill=VIOLET, line=None)
        text(s, 7.75, cy - 0.05, 4.7, 0.35, [[R(t, 11.5, WHITE, True, FONT_SB)]])
        text(s, 7.75, cy + 0.28, 4.7, 0.4, [[R(b, 9.8, GRAY, False, FONT)]], line_spacing=1.0)
    footer(s, 16)


# ----------------------------------------------------------------------------
# Slide 13 — Closing
# ----------------------------------------------------------------------------
def s13():
    s = slide()
    oval(s, -1.5, 3.5, 6, 6, fill=None, line=RGBColor(0x12,0x24,0x3C), lw=1.2)
    oval(s, 9.5, -2.2, 6, 6, fill=None, line=RGBColor(0x12,0x24,0x3C), lw=1.2)

    line(s, 4.3, 2.2, 9.0, 2.2, CYAN, 3)
    text(s, 1.0, 2.35, 11.3, 1.2, [[R("AEGIS", 62, WHITE, True, FONT_SB),
                                    R("-X", 62, CYAN, True, FONT_SB)]], align=PP_ALIGN.CENTER)
    text(s, 1.0, 3.6, 11.3, 0.6,
         [[R("From terabytes to testimony — in hours, not months.", 19, CYAN, True, FONT_SB)]],
         align=PP_ALIGN.CENTER)
    rect(s, 2.9, 4.45, 7.5, 0.8, fill=PANEL2, line=VIOLET, lw=1.5, radius=0.22)
    text(s, 2.9, 4.45, 7.5, 0.8,
         [[R("Every hour saved is a child protected sooner.", 16, WHITE, True, FONT_SB)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    text(s, 1.0, 5.5, 11.3, 0.7,
         [[R("Every deployment makes every other agency stronger —", 15, CYAN, True, FONT_SB)],
          [R("without a single file leaving home.", 15, LGRAY, False, FONT)]],
         align=PP_ALIGN.CENTER, space_after=2)
    text(s, 1.0, 6.5, 11.3, 0.4,
         [[R("HACK-KP 2026", 11.5, GRAY, True, FONT_SB),
           R("   ·   Sovereign · Agentic · Court-Provable Fusion", 11.5, GRAY, False, FONT)]],
         align=PP_ALIGN.CENTER)
    footer(s, 17)


# ----------------------------------------------------------------------------
# Slide A — The Whitespace (lane chart)
# ----------------------------------------------------------------------------
def sA():
    s = slide()
    title_block(s, "The Whitespace", "Everyone owns one lane. The fusion box is empty.", "11")
    text(s, 0.7, 1.92, 11.9, 0.4,
         [[R("Incumbents each dominate a single domain. Nobody owns the ", 12.5, LGRAY, False, FONT),
           R("combination", 12.5, CYAN, True, FONT_SB),
           R(" — sovereign, agentic, court-provable fusion.", 12.5, LGRAY, False, FONT)]])

    MBLUE = RGBColor(0x60, 0xA5, 0xFA)
    lanes = [
        ("Cellebrite", "Devices & mobile extraction", AMBER),
        ("Chainalysis", "Cryptocurrency tracing", GREEN),
        ("Clearview AI", "Face recognition", VIOLET),
        ("Palantir", "Data ontology / fusion UI", MBLUE),
        ("Thorn / PhotoDNA", "Known-hash matching", RED),
    ]
    x0 = 0.7; laneW = 11.9; y0 = 2.55; lh = 0.58; gap = 0.14
    for i, (co, dom, col) in enumerate(lanes):
        cy = y0 + i * (lh + gap)
        rect(s, x0, cy, laneW, lh, fill=PANEL, line=STROKE, lw=1, radius=0.09)
        rect(s, x0, cy, 0.13, lh, fill=col, line=None, rounded=False)
        text(s, x0 + 0.38, cy, 3.5, lh, [[R(co, 13.5, WHITE, True, FONT_SB)]], anchor=MSO_ANCHOR.MIDDLE)
        text(s, x0 + 4.0, cy, 3.6, lh, [[R(dom, 11, GRAY, False, FONT)]], anchor=MSO_ANCHOR.MIDDLE)

    # Vertical AEGIS-X bar crossing every lane (drawn on top)
    vx = 8.75; vy = y0 - 0.03; vw = 1.55
    vh = (lh + gap) * len(lanes) - gap + 0.06
    rect(s, vx, vy, vw, vh, fill=CYAN_D, line=CYAN, lw=2.25, radius=0.06)
    text(s, vx, vy + vh/2 - 0.55, vw, 1.2,
         [[R("AEGIS-X", 16, WHITE, True, FONT_SB)],
          [R("FUSION", 11.5, NAVY, True, FONT_SB)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=1)

    rect(s, 0.7, 6.1, 11.93, 0.72, fill=PANEL2, line=AMBER, lw=1.25, radius=0.12)
    text(s, 1.0, 6.1, 11.4, 0.72,
         [[R("The empty box: ", 11.5, AMBER, True, FONT_SB),
           R("Sovereign + Agentic + Court-Provable Fusion.  ", 11.5, WHITE, True, FONT_SB),
           R("Cellebrite shipped agentic AI in Guardian (2025) — the race is on; the fusion box is still empty.",
             11, GRAY, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 12)


# ----------------------------------------------------------------------------
# Slide B — The Provenance Spine (proof chain)
# ----------------------------------------------------------------------------
def sB():
    s = slide()
    title_block(s, "The Provenance Spine · The Moat",
                "Competitors ask courts for trust. We hand them proof.", "12")
    text(s, 0.7, 1.95, 11.9, 0.4,
         [[R("Every AI claim is compiled into a symbolic derivation where every leaf is a hash-sealed exhibit — "
             "the only court-provable AI reasoning chain in the category.", 12.5, LGRAY, False, FONT)]],
         line_spacing=1.05)

    steps = [
        ("AI Claim", "claim(fact, rule,\nevidence, model, conf.)", CYAN),
        ("Statute Element", "POCSO §14 —\n\"production\"", VIOLET),
        ("Evidence ID", "FILE-2291 ·\nDEV-01 (EXIF match)", AMBER),
        ("SHA-256 Hash", "0x7f3a…c19\nsealed ✓", GREEN),
        ("Enclave Attestation", "NVIDIA CC ·\nmeasurement 0x7f3a…", MBLUE_B()),
    ]
    x0 = 0.7; bw = 2.15; gapx = 0.28; y = 2.85; h = 1.95
    for i, (name, detail, col) in enumerate(steps):
        cx = x0 + i * (bw + gapx)
        rect(s, cx, y, bw, h, fill=PANEL, line=col, lw=1.5, radius=0.07)
        oval(s, cx + bw/2 - 0.27, y + 0.22, 0.54, 0.54, fill=None, line=col, lw=1.75)
        text(s, cx + bw/2 - 0.27, y + 0.22, 0.54, 0.54, [[R(str(i+1), 16, col, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        text(s, cx + 0.1, y + 0.9, bw - 0.2, 0.5, [[R(name, 12.5, WHITE, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, line_spacing=0.95)
        text(s, cx + 0.12, y + 1.38, bw - 0.24, 0.5, [[R(detail, 9.5, GRAY, False, FONT)]],
             align=PP_ALIGN.CENTER, line_spacing=0.95)
        if i < len(steps) - 1:
            ax = cx + bw
            line(s, ax + 0.02, y + h/2, ax + gapx - 0.02, y + h/2, CYAN, 2.25)
            tri = s.shapes.add_shape(MSO_SHAPE.ISOSCELES_TRIANGLE, Inches(ax + gapx - 0.15),
                                     Inches(y + h/2 - 0.08), Inches(0.16), Inches(0.16))
            tri.rotation = 90; tri.fill.solid(); tri.fill.fore_color.rgb = CYAN
            tri.line.fill.background(); tri.shadow.inherit = False

    rect(s, 0.7, 5.5, 11.93, 1.25, fill=PANEL2, line=GREEN, lw=1.5, radius=0.08)
    chip(s, 1.0, 5.72, 2.55, 0.44, "VERIFIER AGENT", GREEN, GREEN)
    text(s, 3.75, 5.5, 8.7, 1.25,
         [[R("No claim reaches an investigator until it is re-derived against raw evidence.", 12.5, WHITE, True, FONT_SB)],
          [R("Every reasoning trace is persisted to the audit log — courts can replay the AI's \"thought process.\"",
             11, GRAY, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=3)
    footer(s, 13)


# ----------------------------------------------------------------------------
# Slide C — Guardian Network (federation)
# ----------------------------------------------------------------------------
def sC():
    import math
    s = slide()
    title_block(s, "Guardian Network · Network Effect",
                "Every deployment makes every other agency stronger", "13")
    chip(s, 0.7, 1.95, 8.6, 0.44,
         "Arrows carry DP-noised gradients only — raw data never leaves the premises", CYAN)

    cx, cy = 6.55, 4.15
    rx, ry = 3.7, 1.45
    nodes = ["State CID", "Cyber Cell", "I4C / CCTNS", "INTERPOL NCB", "Forensic Lab", "Anti-Traffick Unit"]
    node_colors = [AMBER, GREEN, CYAN, VIOLET, MBLUE_B(), RGBColor(0xF4, 0x72, 0xB6)]
    pos = []
    for i in range(len(nodes)):
        a = -math.pi/2 + i * (2*math.pi/len(nodes))
        nx = cx + rx * math.cos(a)
        ny = cy + ry * math.sin(a)
        pos.append((nx, ny))
    # connectors first
    for (nx, ny) in pos:
        line(s, cx, cy, nx, ny, STROKE, 1.5, dash="dash")
    # agency nodes
    nw, nh = 1.85, 0.78
    for (nx, ny), label, col in zip(pos, nodes, node_colors):
        rect(s, nx - nw/2, ny - nh/2, nw, nh, fill=PANEL, line=col, lw=1.5, radius=0.18)
        text(s, nx - nw/2, ny - nh/2, nw, nh, [[R(label, 10.5, WHITE, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    # central aggregator
    aw, ah = 2.5, 1.15
    oval(s, cx - aw/2, cy - ah/2, aw, ah, fill=CYAN_D, line=CYAN, lw=2)
    text(s, cx - aw/2, cy - ah/2, aw, ah,
         [[R("SECURE", 11, WHITE, True, FONT_SB)], [R("AGGREGATOR", 11, WHITE, True, FONT_SB)],
          [R("FL + Diff. Privacy", 9.5, RGBColor(0xC7,0xF2,0xFB), False, FONT)]],
         align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE, space_after=0)

    rect(s, 0.7, 6.35, 11.93, 0.72, fill=PANEL2, line=VIOLET, lw=1.5, radius=0.12)
    text(s, 1.0, 6.35, 11.4, 0.72,
         [[R("Agency #200 benefits from the experience of 199 others.  ", 12, WHITE, True, FONT_SB),
           R("The Flock Safety pattern: data-network effects → 25× revenue multiples.", 11.5, VIOLET, True, FONT_SB)]],
         anchor=MSO_ANCHOR.MIDDLE)
    footer(s, 14)


# ----------------------------------------------------------------------------
# Slide D — Billion-Dollar Math (valuation table)
# ----------------------------------------------------------------------------
def sD():
    from pptx.util import Inches as In
    s = slide()
    title_block(s, "Billion-Dollar Math", "Modest penetration clears $1B", "14")

    rows = [
        ("Scenario", "Agencies × ACV", "ARR", "Multiple", "Valuation"),
        ("Conservative", "200 × $500k", "$100M", "10×", "$1.0B"),
        ("Base", "300 × $500k", "$150M", "12×", "$1.8B"),
        ("Growth · network effect", "500 × $800k", "$400M", "15–25×", "$6B+"),
    ]
    left, top = 0.7, 2.15
    tw, th = 11.93, 2.5
    gt = s.shapes.add_table(len(rows), 5, In(left), In(top), In(tw), In(th)).table
    widths = [3.1, 2.9, 1.9, 1.83, 2.2]
    for c, wv in enumerate(widths):
        gt.columns[c].width = In(wv)
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = gt.cell(r, c)
            cell.margin_left = In(0.16); cell.margin_right = In(0.1)
            cell.margin_top = In(0.04); cell.margin_bottom = In(0.04)
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            growth = (r == 3)
            if r == 0:
                cell.fill.solid(); cell.fill.fore_color.rgb = CYAN_D if c == 4 else PANEL2
            elif c == 4:
                cell.fill.solid(); cell.fill.fore_color.rgb = RGBColor(0x0F, 0x2B, 0x33)
            elif growth:
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL2
            else:
                cell.fill.solid(); cell.fill.fore_color.rgb = PANEL
            tf = cell.text_frame; tf.word_wrap = True
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.CENTER
            run = p.add_run(); run.text = val; run.font.name = FONT
            if r == 0:
                run.font.size = Pt(12); run.font.bold = True
                run.font.color.rgb = NAVY if c == 4 else WHITE
            elif c == 4:
                run.font.size = Pt(17 if growth else 15); run.font.bold = True
                run.font.color.rgb = GREEN if growth else CYAN
            elif c == 0:
                run.font.size = Pt(12.5); run.font.bold = True
                run.font.color.rgb = GREEN if growth else WHITE
            else:
                run.font.size = Pt(12); run.font.color.rgb = LGRAY if growth else GRAY

    # Anchors + TAM band
    rect(s, 0.7, 4.95, 11.93, 1.0, fill=PANEL2, line=CYAN, lw=1.25, radius=0.1)
    text(s, 1.0, 4.95, 2.9, 1.0,
         [[R("TAM", 11, GRAY, False, FONT)], [R("$36B+", 26, CYAN, True, FONT_SB)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=1)
    line(s, 4.1, 5.15, 4.1, 5.75, STROKE, 1)
    text(s, 4.35, 4.95, 8.1, 1.0,
         [[R("PUBLIC-COMP ANCHORS", 9.5, GRAY, True, FONT_SB)],
          [R("Cellebrite ", 11.5, WHITE, True, FONT_SB), R("$480.8M ARR (~6× floor)   ·   ", 11, GRAY, False, FONT),
           R("Magnet ", 11.5, WHITE, True, FONT_SB), R("$1.3B exit (~8–9×)   ·   ", 11, GRAY, False, FONT),
           R("Flock ", 11.5, WHITE, True, FONT_SB), R("$7.5B (~25×, network-effect ceiling)", 11, GRAY, False, FONT)]],
         anchor=MSO_ANCHOR.MIDDLE, space_after=3)

    # GTM strip
    text(s, 0.7, 6.12, 3, 0.3, [[R("GO-TO-MARKET", 10, CYAN, True, FONT_SB)]])
    gtm = [("Kerala cyber-cell pilot", CYAN), ("I4C / CCTNS national", AMBER), ("INTERPOL / ICSE global", GREEN)]
    gx = 0.7; gw = 3.55; ggap = 0.45; gy = 6.42
    for i, (label, col) in enumerate(gtm):
        cxx = gx + i * (gw + ggap)
        rect(s, cxx, gy, gw, 0.5, fill=PANEL, line=col, lw=1.5, radius=0.3)
        text(s, cxx, gy, gw, 0.5, [[R(label, 11.5, WHITE, True, FONT_SB)]],
             align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        if i < len(gtm) - 1:
            ax = cxx + gw
            line(s, ax + 0.05, gy + 0.25, ax + ggap - 0.05, gy + 0.25, col, 2.5)
    footer(s, 15)


def MBLUE_B():
    return RGBColor(0x60, 0xA5, 0xFA)


# ----------------------------------------------------------------------------
for fn in (s1, s2, s3, s4, s5, s6, s7, s8, s9, s10, s11, sA, sB, sC, sD, s12, s13):
    fn()

OUT = os.path.join(os.path.dirname(os.path.abspath(__file__)), "AEGIS-pitch.pptx")
prs.save(OUT)
print("Saved:", OUT)
print("Slides:", len(prs.slides._sldIdLst))
